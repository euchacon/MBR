from flask import Flask, request, jsonify, send_file, render_template
from pptx import Presentation
from pptx.dml.color import RGBColor
import copy, io, os, re, json, requests as req

app = Flask(__name__)
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), 'template.pptx')

JSONBIN_API_KEY = os.environ.get('JSONBIN_API_KEY', '')
JSONBIN_BIN_ID  = os.environ.get('JSONBIN_BIN_ID', '')
JSONBIN_BASE    = 'https://api.jsonbin.io/v3'

# ── JSONBin helpers ───────────────────────────────────────────────────────────

def jb_headers():
    return {'X-Master-Key': JSONBIN_API_KEY, 'Content-Type': 'application/json'}

def load_store():
    """Load the full MBR history dict from JSONBin."""
    if not JSONBIN_BIN_ID:
        return {}
    try:
        r = req.get(f'{JSONBIN_BASE}/b/{JSONBIN_BIN_ID}/latest', headers=jb_headers(), timeout=8)
        if r.status_code == 200:
            return r.json().get('record', {})
    except:
        pass
    return {}

def save_store(store):
    """Overwrite the JSONBin bin with updated store."""
    if not JSONBIN_BIN_ID:
        return False
    try:
        r = req.put(f'{JSONBIN_BASE}/b/{JSONBIN_BIN_ID}', headers=jb_headers(),
                    data=json.dumps(store), timeout=8)
        return r.status_code == 200
    except:
        return False

# ── PPTX helpers ──────────────────────────────────────────────────────────────

def set_text_frame(tf, lines):
    if not lines:
        lines = ['']
    from pptx.oxml.ns import qn
    from lxml import etree
    ref_para = tf.paragraphs[0] if tf.paragraphs else None
    ref_run  = ref_para.runs[0] if (ref_para and ref_para.runs) else None
    txBody   = tf._txBody
    existing = txBody.findall(qn('a:p'))
    for p in existing[1:]:
        txBody.remove(p)
    first_p = existing[0]
    for r in first_p.findall(qn('a:r')):
        first_p.remove(r)
    r_elem = copy.deepcopy(ref_run._r) if ref_run else etree.SubElement(first_p, qn('a:r'))
    t_elem = r_elem.find(qn('a:t'))
    if t_elem is None:
        t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = lines[0]
    first_p.append(r_elem)
    for line in lines[1:]:
        new_p = copy.deepcopy(first_p)
        for r in new_p.findall(qn('a:r')):
            new_p.remove(r)
        new_r = copy.deepcopy(r_elem)
        new_t = new_r.find(qn('a:t'))
        if new_t is None:
            new_t = etree.SubElement(new_r, qn('a:t'))
        new_t.text = line
        new_p.append(new_r)
        txBody.append(new_p)

def replace_shape_text(slide, shape_name, new_text):
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            lines = new_text.split('\n') if new_text else ['']
            set_text_frame(shape.text_frame, lines)
            return True
    return False

def fill_slide_table(slide, table_index, rows_data, col_keys):
    from pptx.oxml.ns import qn
    from lxml import etree
    tables = [s for s in slide.shapes if s.has_table]
    if table_index >= len(tables):
        return
    tbl = tables[table_index].table
    for ri, row_data in enumerate(rows_data):
        tbl_row = ri + 1
        if tbl_row >= len(tbl.rows):
            break
        for ci, col_key in enumerate(col_keys):
            if ci >= len(tbl.columns):
                break
            cell   = tbl.cell(tbl_row, ci)
            tf     = cell.text_frame
            txBody = tf._txBody
            paras  = txBody.findall(qn('a:p'))
            for p in paras[1:]:
                txBody.remove(p)
            first_p = paras[0]
            for r in first_p.findall(qn('a:r')):
                first_p.remove(r)
            r_elem = etree.SubElement(first_p, qn('a:r'))
            t_elem = etree.SubElement(r_elem, qn('a:t'))
            t_elem.text = str(row_data.get(col_key, ''))

# ── Routes ────────────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/save', methods=['POST'])
def save_mbr():
    """Save MBR data for a given period to JSONBin."""
    d      = request.json
    period = d.get('period', '')
    if not period:
        return jsonify({'ok': False, 'error': 'No period provided'}), 400
    store          = load_store()
    store[period]  = d
    ok             = save_store(store)
    return jsonify({'ok': ok})


@app.route('/load/<path:period>', methods=['GET'])
def load_mbr(period):
    """Load MBR data for a given period from JSONBin."""
    store = load_store()
    data  = store.get(period)
    if data is None:
        return jsonify({'ok': False, 'error': 'Not found'}), 404
    return jsonify({'ok': True, 'data': data})


@app.route('/history', methods=['GET'])
def history():
    """Return list of saved periods."""
    store   = load_store()
    periods = list(store.keys())
    return jsonify({'ok': True, 'periods': periods})


@app.route('/generate', methods=['POST'])
def generate():
    d   = request.json
    prs = Presentation(TEMPLATE_PATH)
    s   = prs.slides

    # Slide 1 — Cover
    replace_shape_text(s[0], 'Text 6',  d.get('dealer', ''))
    replace_shape_text(s[0], 'Text 8',  d.get('country', ''))
    replace_shape_text(s[0], 'Text 10', d.get('period', ''))
    replace_shape_text(s[0], 'Text 12', d.get('prepared_by', 'Eugenio Chacon'))

    # Slide 2 — Snapshot
    replace_shape_text(s[1], 'Text 4',  d.get('active_accounts', '[#]'))
    replace_shape_text(s[1], 'Text 8',  d.get('total_backlog', '$[X,XXX]'))
    replace_shape_text(s[1], 'Text 12', d.get('units_invoiced', '[#]'))
    replace_shape_text(s[1], 'Text 16', d.get('revenue_invoiced', '$[X,XXX]'))
    replace_shape_text(s[1], 'Text 22', d.get('key_movements', ''))
    replace_shape_text(s[1], 'Text 27', d.get('lzb_needs', ''))
    replace_shape_text(s[1], 'Text 30', d.get('blockers', ''))

    # Slide 3 — Account Performance
    fill_slide_table(s[2], 0, d.get('accounts', []),
        ['account','doors','slots','units_sold','revenue','floor_inv','backlog','rotation','vs_target','status'])
    replace_shape_text(s[2], 'Text 148', d.get('perf_notes', ''))

    # Slide 4 — Inventory
    replace_shape_text(s[3], 'Text 4',  d.get('units_on_floor', '[#]'))
    replace_shape_text(s[3], 'Text 8',  d.get('units_in_transit', '[#]'))
    replace_shape_text(s[3], 'Text 12', d.get('backlog_value', '$[X,XXX]'))
    replace_shape_text(s[3], 'Text 16', d.get('weeks_cover', '[#]') + ' wks')
    fill_slide_table(s[3], 0, d.get('inventory', []),
        ['account','sku','on_floor','sold_mtd','in_transit','reorder','next_action','target_stock','days_cover'])
    replace_shape_text(s[3], 'Text 114', d.get('inv_issues', ''))

    # Slide 5 — Commercial Narrative
    replace_shape_text(s[4], 'Text 5',  d.get('what_worked', ''))
    replace_shape_text(s[4], 'Text 9',  d.get('what_didnt', ''))
    replace_shape_text(s[4], 'Text 13', d.get('cust_feedback', ''))
    replace_shape_text(s[4], 'Text 17', d.get('market_obs', ''))

    # Slide 6 — Pipeline
    fill_slide_table(s[5], 0, d.get('pipeline', []),
        ['account','opportunity','est_units','est_usd','probability','expected_close','next_action','owner'])
    replace_shape_text(s[5], 'Text 88', d.get('orders_to_place', ''))
    replace_shape_text(s[5], 'Text 91', d.get('activities_planned', ''))

    # Slide 7 — New Customer Deployment
    fill_slide_table(s[6], 0, d.get('new_customers', []),
        ['market','dealer','customer','doors','total_slots','launch_date','status','next_step'])
    replace_shape_text(s[6], 'Text 71', d.get('onboarding_notes', ''))

    # Slide 8 — Needs & Escalations
    replace_shape_text(s[7], 'Text 5',  d.get('need_marketing', ''))
    replace_shape_text(s[7], 'Text 9',  d.get('need_pricing', ''))
    replace_shape_text(s[7], 'Text 13', d.get('need_product', ''))
    replace_shape_text(s[7], 'Text 17', d.get('need_training', ''))
    replace_shape_text(s[7], 'Text 21', d.get('need_logistics', ''))
    replace_shape_text(s[7], 'Text 25', d.get('need_strategic', ''))

    # Slide 9 — Marketing
    replace_shape_text(s[8], 'Text 4',  d.get('mkt_events', ''))
    replace_shape_text(s[8], 'Text 7',  d.get('mkt_digital', ''))
    replace_shape_text(s[8], 'Text 10', d.get('mkt_instore', ''))
    replace_shape_text(s[8], 'Text 15', d.get('mkt_asks', ''))

    # Slide 10 — Closing
    replace_shape_text(s[9], 'Text 2',
        f'Next review:  {d.get("next_review","[Date]")}  -  {d.get("next_format","video call")}')
    replace_shape_text(s[9], 'Text 3',
        f'Questions or follow-ups:  {d.get("dealer_email","")}  |  Jeff:  {d.get("jeff_email","jeff.lillich@la-z-boy.com")}  |  LATAM Team:  {d.get("latam_email","eugenio.chacon@gmail.com")}')
    replace_shape_text(s[9], 'Text 4',
        f'La-Z-Boy LATAM  -  Confidential. Internal Use Only.  -  {d.get("period","[Month YYYY]")}')

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    period_safe = re.sub(r'[^a-zA-Z0-9_-]', '_', d.get('period', 'MBR'))
    return send_file(buf, as_attachment=True, download_name=f'LZB_LATAM_MBR_{period_safe}.pptx',
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


if __name__ == '__main__':
    app.run(debug=True, port=5050)
